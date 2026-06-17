import io
import os
import tempfile
from functools import wraps
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from flask import Flask, render_template, request, session, redirect, g
from llm_client import predict_justification_score_api, predict_justification_evaluation_api, generate_mcqs_from_notes_api
import mysql.connector
from config import DB_CONFIG

def get_db_connection():
    return mysql.connector.connect(**DB_CONFIG)

app = Flask(__name__)
app.secret_key = 'secret'
# -------------------------------------------------------
# EduSense Grading Engine Configuration
# -------------------------------------------------------

GRADING_ENGINE = "llm_api"  # Options: "llm_api" or "cosine"

# -------------------------------------------------------
# EduSense User Context
# -------------------------------------------------------

@app.before_request
def load_logged_in_user():
    """
    Loads the signed-in user name for navbar greeting.
    This keeps the LMS interface personal and professional.
    """
    g.current_user = None

    user_role = session.get("user_role")
    user_id = session.get("user_id")

    if not user_role or not user_id:
        return

    try:
        conn = get_db_connection()
        cursor = conn.cursor(dictionary=True)

        if user_role == "student":
            cursor.execute("""
                SELECT id, name, username, roll_number
                FROM students
                WHERE id = %s
            """, (user_id,))
            user = cursor.fetchone()

            if user:
                g.current_user = {
                    "id": user["id"],
                    "name": user["name"] or user["username"] or "Student",
                    "username": user["username"],
                    "role": "student",
                    "role_label": "Student",
                    "roll_number": user.get("roll_number")
                }

        elif user_role in ["administrator", "teacher"]:
            cursor.execute("""
                SELECT id, name, username, role
                FROM admins
                WHERE id = %s
            """, (user_id,))
            user = cursor.fetchone()

            if user:
                role_label = "Administrator" if user_role == "administrator" else "Teacher"
                g.current_user = {
                    "id": user["id"],
                    "name": user.get("name") or user["username"] or role_label,
                    "username": user["username"],
                    "role": user_role,
                    "role_label": role_label
                }

        cursor.close()
        conn.close()

    except Exception as e:
        print("Could not load current user:", e)
        g.current_user = {
            "id": user_id,
            "name": "EduSense User",
            "username": "",
            "role": user_role,
            "role_label": user_role.title()
        }


@app.context_processor
def inject_current_user():
    return {
        "current_user": g.get("current_user", None)
    }

# -------------------------------------------------------
# EduSense Access Control Helpers
# -------------------------------------------------------

def redirect_to_role_dashboard():
    """
    Sends already logged-in users to the correct dashboard.
    This prevents students/admins/teachers from opening the wrong area.
    """
    role = session.get("user_role")

    if role in ["administrator", "teacher"]:
        return redirect("/admin/dashboard")

    if role == "student":
        return redirect("/student/dashboard")

    return redirect("/login")


def role_required(*allowed_roles):
    """
    Reusable route protection decorator.
    Example:
        @role_required("student")
        @role_required("administrator", "teacher")
    """
    def decorator(view_func):
        @wraps(view_func)
        def wrapped_view(*args, **kwargs):
            user_role = session.get("user_role")
            user_id = session.get("user_id")

            if not user_role or not user_id:
                return redirect("/login")

            if user_role not in allowed_roles:
                return redirect_to_role_dashboard()

            return view_func(*args, **kwargs)

        return wrapped_view

    return decorator


@app.after_request
def add_no_cache_headers(response):
    """
    Prevents browser back-button from showing protected pages after logout.
    """
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    response.headers["Pragma"] = "no-cache"
    response.headers["Expires"] = "0"
    return response

def build_options_text(mcq):
    """
    Build MCQ options text for LLM prompt.
    """
    return (
        f"A) {mcq.get('option_a', '')} "
        f"B) {mcq.get('option_b', '')} "
        f"C) {mcq.get('option_c', '')} "
        f"D) {mcq.get('option_d', '')}"
    )


def cosine_justification_score(model_justification, student_justification):
    """
    Existing TF-IDF cosine fallback.
    Returns score from 0 to 100.
    """
    if not student_justification or not student_justification.strip():
        return 0.0

    try:
        vectorizer = TfidfVectorizer()
        tfidf = vectorizer.fit_transform([model_justification or "", student_justification])
        similarity = cosine_similarity(tfidf[0:1], tfidf[1:2])[0][0]
        return float(similarity * 100)
    except Exception as e:
        print("Cosine scoring failed:", str(e))
        return 0.0


def get_justification_score(mcq, student_justification):
    """
    Main justification scoring function.
    First tries LLM API.
    If LLM API fails, falls back to cosine similarity.
    Returns:
        just_percent, score_source
    """

    if not student_justification or not student_justification.strip():
        return 0.0, "empty"

    if GRADING_ENGINE == "llm_api":
        options_text = build_options_text(mcq)

        llm_score = predict_justification_score_api(
            question=mcq.get("question", ""),
            options=options_text,
            model_justification=mcq.get("model_justification", ""),
            student_justification=student_justification
        )

        if llm_score is not None:
            return float(llm_score), "llm_api"

        print("LLM API unavailable. Falling back to cosine similarity.")

    cosine_score = cosine_justification_score(
        mcq.get("model_justification", ""),
        student_justification
    )

    return cosine_score, "cosine"

def get_justification_evaluation(mcq, student_justification):
    """
    Returns:
        just_percent, score_source, llm_feedback
    """

    if not student_justification or not student_justification.strip():
        return 0.0, "empty", "No written justification was submitted."

    if GRADING_ENGINE == "llm_api":
        try:
            options_text = build_options_text(mcq)

            evaluation = predict_justification_evaluation_api(
                question=mcq.get("question", ""),
                options=options_text,
                model_justification=mcq.get("model_justification", ""),
                student_justification=student_justification
            )

            score = float(evaluation.get("score", 0.0))
            feedback = evaluation.get("feedback") or "The response was evaluated by the AI scoring model."

            return score, "llm_api", feedback

        except Exception as e:
            print("LLM evaluation failed:", str(e))

    cosine_score = cosine_justification_score(
        mcq.get("model_justification", ""),
        student_justification
    )

    if cosine_score >= 80:
        feedback = "The student explanation closely matches the expected reasoning."
    elif cosine_score >= 50:
        feedback = "The student explanation shows partial understanding but needs more accurate detail."
    elif cosine_score > 0:
        feedback = "The student explanation has limited similarity with the expected reasoning."
    else:
        feedback = "The submitted explanation does not provide enough relevant reasoning."

    return cosine_score, "cosine", feedback

def extract_text_from_uploaded_file(uploaded_file):
    """
    Extract text from uploaded TXT, MD, PDF, DOCX, or PPTX file.
    Returns extracted text as string.
    """

    if not uploaded_file or uploaded_file.filename == "":
        return ""

    filename = uploaded_file.filename.lower()
    file_bytes = uploaded_file.read()

    if not file_bytes:
        return ""

    try:
        if filename.endswith(".txt") or filename.endswith(".md"):
            return file_bytes.decode("utf-8", errors="ignore").strip()

        if filename.endswith(".pdf"):
            reader = PdfReader(io.BytesIO(file_bytes))
            text_parts = []

            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

            return "\n".join(text_parts).strip()

        if filename.endswith(".docx"):
            document = Document(io.BytesIO(file_bytes))
            paragraphs = []

            for paragraph in document.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text.strip())

            return "\n".join(paragraphs).strip()

        if filename.endswith(".pptx"):
            presentation = Presentation(io.BytesIO(file_bytes))
            slide_texts = []

            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

            return "\n".join(slide_texts).strip()

        return ""

    except Exception as e:
        print("File text extraction failed:", str(e))
        return ""


def is_allowed_notes_file(filename):
    """
    Checks whether uploaded file extension is supported.
    """
    if not filename:
        return False

    filename = filename.lower()

    return (
        filename.endswith(".txt") or
        filename.endswith(".md") or
        filename.endswith(".pdf") or
        filename.endswith(".docx") or
        filename.endswith(".pptx")
    )

@app.route('/')
def index():
    return redirect_to_role_dashboard()

@app.route('/login', methods=['GET', 'POST'])
def login_page():
    if session.get("user_role") and session.get("user_id"):
        if session.get("user_role") in ["administrator", "teacher"]:
            return redirect("/admin/dashboard")
        if session.get("user_role") == "student":
            return redirect("/student/dashboard")

    error = None

    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        password = request.form.get('password', '').strip()

        if not username or not password:
            error = "Please enter both username and password."
            return render_template('login.html', error=error)

        conn = get_db_connection()
        cursor = conn.cursor(dictionary=True)

        cursor.execute("""
            SELECT id, role, name, username
            FROM admins
            WHERE username = %s AND password = %s
        """, (username, password))
        admin = cursor.fetchone()

        if admin:
            session.clear()
            session['user_role'] = admin['role']
            session['user_id'] = admin['id']
            session['display_name'] = admin.get('name') or admin.get('username') or 'Administrator'

            cursor.close()
            conn.close()
            return redirect('/admin/dashboard')

        cursor.execute("""
            SELECT id, name, username
            FROM students
            WHERE username = %s AND password = %s
        """, (username, password))
        student = cursor.fetchone()

        if student:
            session.clear()
            session['user_role'] = 'student'
            session['user_id'] = student['id']
            session['display_name'] = student.get('name') or student.get('username') or 'Student'

            cursor.close()
            conn.close()
            return redirect('/student/dashboard')

        cursor.close()
        conn.close()
        error = "Invalid username or password."

    return render_template('login.html', error=error)

@app.route('/student/dashboard')
def student_dashboard():
    if session.get('user_role') != 'student':
        return redirect('/login')

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT 
            s.id,
            s.name,
            s.roll_number,
            s.username,
            c.name AS class_name
        FROM students s
        LEFT JOIN classes c ON s.class_id = c.id
        WHERE s.id = %s
    """, (session['user_id'],))
    student = cursor.fetchone() or {}

    cursor.execute("""
        SELECT COUNT(*) AS total_courses
        FROM books b
        JOIN classes c ON b.class_id = c.id
        JOIN students s ON s.class_id = c.id
        WHERE s.id = %s AND b.is_published = TRUE
    """, (session['user_id'],))
    total_courses = cursor.fetchone()["total_courses"] or 0

    cursor.execute("""
        SELECT COUNT(DISTINCT m.book_id) AS attempted_courses
        FROM results r
        JOIN mcqs m ON r.mcq_id = m.id
        WHERE r.student_id = %s
    """, (session['user_id'],))
    attempted_courses = cursor.fetchone()["attempted_courses"] or 0

    cursor.execute("""
        SELECT COUNT(*) AS answered_items
        FROM results
        WHERE student_id = %s
    """, (session['user_id'],))
    answered_items = cursor.fetchone()["answered_items"] or 0

    cursor.execute("""
        SELECT ROUND(AVG(course_percentage), 2) AS average_percentage
        FROM (
            SELECT SUM(r.marks) / COUNT(r.id) * 100 AS course_percentage
            FROM results r
            JOIN mcqs m ON r.mcq_id = m.id
            WHERE r.student_id = %s
            GROUP BY m.book_id
        ) AS course_scores
    """, (session['user_id'],))
    average_row = cursor.fetchone()
    average_percentage = average_row["average_percentage"] if average_row and average_row["average_percentage"] is not None else 0

    stats = {
        "total_courses": total_courses,
        "attempted_courses": attempted_courses,
        "pending_courses": max(total_courses - attempted_courses, 0),
        "answered_items": answered_items,
        "average_percentage": average_percentage
    }

    cursor.close()
    conn.close()

    return render_template('student/dashboard.html', student=student, stats=stats)

@app.route('/admin/dashboard')
def admin_dashboard():
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT id, name, username, role
        FROM admins
        WHERE id = %s
    """, (session['user_id'],))
    user = cursor.fetchone() or {}

    cursor.execute("SELECT COUNT(*) AS total FROM classes")
    total_classes = cursor.fetchone()["total"] or 0

    cursor.execute("SELECT COUNT(*) AS total FROM books")
    total_courses = cursor.fetchone()["total"] or 0

    cursor.execute("SELECT COUNT(*) AS total FROM books WHERE is_published = TRUE")
    published_courses = cursor.fetchone()["total"] or 0

    cursor.execute("SELECT COUNT(*) AS total FROM students")
    total_students = cursor.fetchone()["total"] or 0

    cursor.execute("SELECT COUNT(*) AS total FROM admins WHERE role = 'teacher'")
    total_teachers = cursor.fetchone()["total"] or 0

    cursor.execute("SELECT COUNT(*) AS total FROM mcqs")
    total_mcqs = cursor.fetchone()["total"] or 0

    cursor.execute("SELECT COUNT(*) AS total FROM results")
    total_results = cursor.fetchone()["total"] or 0

    cursor.execute("""
        SELECT ROUND(AVG(student_percentage), 2) AS average_percentage
        FROM (
            SELECT SUM(r.marks) / COUNT(r.id) * 100 AS student_percentage
            FROM results r
            GROUP BY r.student_id
        ) AS student_scores
    """)
    average_row = cursor.fetchone()
    average_percentage = (
        average_row["average_percentage"]
        if average_row and average_row["average_percentage"] is not None
        else 0
    )

    stats = {
        "total_classes": total_classes,
        "total_courses": total_courses,
        "published_courses": published_courses,
        "total_students": total_students,
        "total_teachers": total_teachers,
        "total_mcqs": total_mcqs,
        "total_results": total_results,
        "average_percentage": average_percentage
    }

    cursor.close()
    conn.close()

    if session['user_role'] == 'administrator':
        return render_template(
            'admin/administrator_dashboard.html',
            user=user,
            stats=stats
        )

    return render_template(
        'admin/teacher_dashboard.html',
        teacher=user,
        stats=stats
    )
    
@app.route('/logout')
def logout():
    session.clear()
    return redirect('/login')

@app.route('/student/profile')
def student_profile():
    if session.get('user_role') != 'student':
        return redirect('/login')
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    cursor.execute("SELECT s.name, s.roll_number, c.name as class_name FROM students s JOIN classes c ON s.class_id = c.id WHERE s.id = %s", (session['user_id'],))
    student = cursor.fetchone()
    conn.close()
    return render_template('student/profile.html', student=student or {})

@app.route('/student/books')
def student_books():
    if session.get('user_role') != 'student':
        return redirect('/login')

    student_id = session['user_id']

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT 
            b.id,
            b.title,
            c.name AS class_name,

            (
                SELECT COUNT(DISTINCT a.id)
                FROM assessments a
                JOIN mcqs m ON m.assessment_id = a.id
                WHERE a.book_id = b.id
                  AND a.is_published = 1
                  AND NOT EXISTS (
                      SELECT 1
                      FROM results r
                      WHERE r.student_id = %s
                        AND r.assessment_id = a.id
                  )
            ) AS pending_count,

            (
                SELECT COUNT(DISTINCT r.assessment_id)
                FROM results r
                JOIN assessments a2 ON r.assessment_id = a2.id
                WHERE r.student_id = %s
                  AND a2.book_id = b.id
                  AND r.assessment_id IS NOT NULL
            ) AS attempted_count

        FROM books b
        JOIN classes c ON b.class_id = c.id
        JOIN students s ON s.class_id = c.id
        WHERE s.id = %s
          AND b.is_published = TRUE
        ORDER BY b.title
    """, (student_id, student_id, student_id))

    books = cursor.fetchall()

    cursor.close()
    conn.close()

    return render_template('student/books.html', books=books)

@app.route('/student/quiz/<int:book_id>/date/<attempt_date>')
def student_quiz_by_date(book_id, attempt_date):
    if session.get('user_role') != 'student':
        return redirect('/login')

    student_id = session['user_id']

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT 
            b.id,
            b.title,
            c.name AS class_name
        FROM books b
        JOIN classes c ON b.class_id = c.id
        JOIN students s ON s.class_id = c.id
        WHERE b.id = %s
          AND s.id = %s
          AND b.is_published = TRUE
    """, (book_id, student_id))
    book = cursor.fetchone()

    if not book:
        cursor.close()
        conn.close()
        return "Quiz not found or not published"

    cursor.execute("""
        SELECT DISTINCT m.*
        FROM mcqs m
        JOIN results r ON r.mcq_id = m.id
        WHERE r.student_id = %s
          AND m.book_id = %s
          AND DATE_FORMAT(r.submitted_at, '%Y-%m-%d') = %s
        ORDER BY m.id
    """, (student_id, book_id, attempt_date))
    mcqs = cursor.fetchall()

    cursor.close()
    conn.close()

    if not mcqs:
        return "No questions found for this quiz date"

    mcq_ids = ",".join(str(mcq["id"]) for mcq in mcqs)

    return render_template(
        'student/quiz.html',
        quiz_name=f"{book['title']} Quiz - {attempt_date}",
        mcqs=mcqs,
        book_id=book_id,
        mcq_ids=mcq_ids,
        quiz_context="reattempt"
    )


@app.route('/student/quiz/<int:book_id>')
def student_quiz(book_id):
    if session.get('user_role') != 'student':
        return redirect('/login')

    student_id = session['user_id']

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT 
            b.id,
            b.title,
            c.name AS class_name
        FROM books b
        JOIN classes c ON b.class_id = c.id
        JOIN students s ON s.class_id = c.id
        WHERE b.id = %s
          AND s.id = %s
          AND b.is_published = TRUE
    """, (book_id, student_id))

    book = cursor.fetchone()

    if not book:
        cursor.close()
        conn.close()
        return "Course not found or not available"

    cursor.execute("""
        SELECT 
            a.id,
            a.title,
            DATE_FORMAT(a.created_at, '%Y-%m-%d') AS created_date,
            COUNT(m.id) AS total_mcqs
        FROM assessments a
        JOIN mcqs m ON m.assessment_id = a.id
        WHERE a.book_id = %s
          AND a.is_published = 1
          AND NOT EXISTS (
              SELECT 1
              FROM results r
              WHERE r.student_id = %s
                AND r.assessment_id = a.id
          )
        GROUP BY a.id, a.title, a.created_at
        HAVING total_mcqs > 0
        ORDER BY a.created_at DESC, a.id DESC
    """, (book_id, student_id))

    assessments = cursor.fetchall()

    cursor.close()
    conn.close()

    return render_template(
        'student/quiz_list.html',
        book=book,
        assessments=assessments
    )

@app.route('/student/assessment/<int:assessment_id>/attempt')
def student_assessment_attempt(assessment_id):
    if session.get('user_role') != 'student':
        return redirect('/login')

    student_id = session['user_id']

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT 
            a.id AS assessment_id,
            a.title AS assessment_title,
            b.id AS book_id,
            b.title AS course_name,
            c.name AS class_name
        FROM assessments a
        JOIN books b ON a.book_id = b.id
        JOIN classes c ON b.class_id = c.id
        JOIN students s ON s.class_id = c.id
        WHERE a.id = %s
          AND s.id = %s
          AND a.is_published = 1
          AND b.is_published = TRUE
    """, (assessment_id, student_id))
    assessment = cursor.fetchone()

    if not assessment:
        cursor.close()
        conn.close()
        return "Assessment not found or not available"

    cursor.execute("""
        SELECT COUNT(*) AS total
        FROM results r
        JOIN mcqs m ON r.mcq_id = m.id
        WHERE r.student_id = %s
          AND m.assessment_id = %s
    """, (student_id, assessment_id))
    already_attempted = cursor.fetchone()["total"] or 0

    if already_attempted > 0:
        cursor.close()
        conn.close()
        return redirect(f"/student/results/detail/{assessment['book_id']}")

    cursor.execute("""
        SELECT 
            id,
            question,
            option_a,
            option_b,
            option_c,
            option_d
        FROM mcqs
        WHERE assessment_id = %s
        ORDER BY id
    """, (assessment_id,))
    mcqs = cursor.fetchall()

    cursor.close()
    conn.close()

    if not mcqs:
        return "No MCQs available in this assessment"

    mcq_ids = ",".join(str(mcq["id"]) for mcq in mcqs)

    return render_template(
        'student/quiz.html',
        quiz_name=assessment["assessment_title"],
        course_name=assessment["course_name"],
        class_name=assessment["class_name"],
        assessment_id=assessment_id,
        book_id=assessment["book_id"],
        mcqs=mcqs,
        mcq_ids=mcq_ids,
        quiz_context="saved_assessment"
    )

@app.route('/student/quiz/submit', methods=['POST'])
def student_quiz_submit():
    if session.get('user_role') != 'student':
        return redirect('/login')

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    student_id = session['user_id']
    book_id = request.form.get('book_id')
    assessment_id = request.form.get('assessment_id') or None
    mcq_ids_raw = request.form.get('mcq_ids', '').strip()

    if not book_id:
        cursor.close()
        conn.close()
        return 'Error: Course ID missing'

    # Convert assessment_id safely
    try:
        assessment_id_int = int(assessment_id) if assessment_id else None
    except Exception:
        assessment_id_int = None

    # If this is a saved assessment and the student already attempted it,
    # do not insert duplicate results.
    if assessment_id_int:
        cursor.execute("""
            SELECT COUNT(*) AS total
            FROM results
            WHERE student_id = %s
              AND assessment_id = %s
        """, (student_id, assessment_id_int))

        already_attempted = cursor.fetchone()["total"] or 0

        if already_attempted > 0:
            cursor.close()
            conn.close()
            return redirect(f"/student/results/detail/{book_id}")

    # Load only the MCQs sent from the form.
    # This prevents the old problem where all course MCQs were submitted together.
    mcq_ids = []

    if mcq_ids_raw:
        mcq_ids = [
            int(item.strip())
            for item in mcq_ids_raw.split(',')
            if item.strip().isdigit()
        ]

    if mcq_ids:
        placeholders = ",".join(["%s"] * len(mcq_ids))

        cursor.execute(f"""
            SELECT 
                id,
                question,
                option_a,
                option_b,
                option_c,
                option_d,
                correct_option,
                model_justification,
                book_id,
                assessment_id
            FROM mcqs
            WHERE id IN ({placeholders})
            ORDER BY id
        """, tuple(mcq_ids))

    elif assessment_id_int:
        cursor.execute("""
            SELECT 
                id,
                question,
                option_a,
                option_b,
                option_c,
                option_d,
                correct_option,
                model_justification,
                book_id,
                assessment_id
            FROM mcqs
            WHERE assessment_id = %s
            ORDER BY id
        """, (assessment_id_int,))

    else:
        # Legacy fallback only. This should rarely run after assessment structure is active.
        cursor.execute("""
            SELECT 
                id,
                question,
                option_a,
                option_b,
                option_c,
                option_d,
                correct_option,
                model_justification,
                book_id,
                assessment_id
            FROM mcqs
            WHERE book_id = %s
            ORDER BY id
        """, (book_id,))

    all_mcqs = cursor.fetchall()

    if not all_mcqs:
        cursor.close()
        conn.close()
        return 'Error: No MCQs found for this assessment'

    total_questions = len(all_mcqs)
    current_results = []
    total_marks = 0.0

    for mcq in all_mcqs:
        mcq_id = mcq['id']
        selected = request.form.get(f'option_{mcq_id}')
        justification = request.form.get(f'justification_{mcq_id}', '').strip()

        option_correct = False
        option_marks = 0.0
        just_percent = 0.0
        just_marks = 0.0
        marks = 0.0
        score_source = "empty"
        feedback = "Skipped (0/1)"

        llm_feedback = "No written justification was submitted."
        if selected or justification:
            option_correct = (selected == mcq['correct_option'])
            option_marks = 0.4 if option_correct else 0.0

            if justification:
                try:
                    just_percent, score_source, llm_feedback = get_justification_evaluation(mcq, justification)
                except Exception as e:
                    print("Justification scoring failed:", str(e))
                    just_percent = 0.0
                    score_source = "error"

            just_marks = (float(just_percent) / 100.0) * 0.6
            marks = option_marks + just_marks

            if score_source == "llm_api":
                justification_label = "LLM Justification"
            elif score_source == "cosine":
                justification_label = "Cosine Fallback Justification"
            elif score_source == "empty":
                justification_label = "No Justification"
            elif score_source == "error":
                justification_label = "Scoring Error"
            else:
                justification_label = "Justification"

            feedback_parts = [
                f"{'Correct' if option_correct else 'Wrong'} option ({option_marks:.1f}/0.4)",
                f"{justification_label}: {just_percent:.1f}% ({just_marks:.1f}/0.6)"
            ]

            feedback = '; '.join(feedback_parts)

        total_marks += marks

        current_results.append({
            'question': mcq['question'],
            'option_a': mcq.get('option_a'),
            'option_b': mcq.get('option_b'),
            'option_c': mcq.get('option_c'),
            'option_d': mcq.get('option_d'),
            'correct_option': mcq.get('correct_option'),
            'selected': selected or 'Not selected',
            'student_justification': justification or 'No justification submitted.',
            'model_justification': mcq.get('model_justification') or '',
            'marks': round(marks, 2),
            'option_marks': round(option_marks, 2),
            'justification_marks': round(just_marks, 2),
            'justification_percent': round(float(just_percent), 2),
            'feedback': feedback
        })

        # Save assessment_id with each result so this quiz disappears from pending assessments.
        cursor.execute("""
            INSERT INTO results 
            (student_id, mcq_id, assessment_id, selected_option, justification, marks, feedback, llm_feedback)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
        """, (
            student_id,
            mcq_id,
            assessment_id_int,
            selected,
            justification,
            round(marks, 4),
            feedback,
            llm_feedback
        ))

    total_out_of = total_questions * 1.0

    conn.commit()
    cursor.close()
    conn.close()

    return render_template(
        'student/current_result.html',
        results=current_results,
        total_marks=round(total_marks, 2),
        total_out_of=total_out_of,
        num_questions=total_questions
    )

@app.route('/student/results')
def student_results():
    if session.get('user_role') != 'student':
        return redirect('/login')

    student_id = session['user_id']

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT 
            b.id AS book_id,
            b.title AS course_name,
            c.name AS class_name,
            COUNT(DISTINCT r.assessment_id) AS attempted_count,
            ROUND(SUM(r.marks), 2) AS obtained_marks,
            COUNT(r.id) AS total_questions,
            ROUND((SUM(r.marks) / COUNT(r.id)) * 100, 2) AS overall_percentage,
            DATE_FORMAT(MAX(r.submitted_at), '%Y-%m-%d') AS last_attempt_date
        FROM results r
        JOIN assessments a ON r.assessment_id = a.id
        JOIN books b ON a.book_id = b.id
        JOIN classes c ON b.class_id = c.id
        WHERE r.student_id = %s
          AND r.assessment_id IS NOT NULL
        GROUP BY b.id, b.title, c.name
        ORDER BY MAX(r.submitted_at) DESC, b.title
    """, (student_id,))

    courses = cursor.fetchall()

    cursor.close()
    conn.close()

    return render_template('student/results.html', courses=courses)


# Admin Routes - Administrator & Teacher access
@app.route('/admin/classes')
def admin_classes():
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')
    # Teachers see all classes (read-only)
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    cursor.execute("SELECT * FROM classes ORDER BY name")
    classes = cursor.fetchall()
    conn.close()
    return render_template('admin/classes.html', classes=classes)

@app.route('/admin/classes/create', methods=['POST'])
def admin_classes_create():
    if session.get('user_role') != 'administrator':
        return redirect('/login')
    name = request.form['name']
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("INSERT INTO classes (name, created_by) VALUES (%s, %s)", (name, session['user_id']))
    conn.commit()
    conn.close()
    return redirect('/admin/classes')

@app.route('/admin/books/<int:class_id>')
def admin_books(class_id):
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    # Remove created_by check for teacher
    cursor.execute("SELECT name FROM classes WHERE id = %s", (class_id,))
    cls = cursor.fetchone()
    if not cls:
        conn.close()
        return 'Class not found'
    cursor.execute("SELECT * FROM books WHERE class_id = %s ORDER BY title", (class_id,))
    books = cursor.fetchall()
    conn.close()
    return render_template('admin/books.html', class_name=cls['name'], books=books, class_id=class_id)

@app.route('/admin/books/create/<int:class_id>', methods=['POST'])
def admin_books_create(class_id):
    if session.get('user_role') != 'administrator':
        return redirect('/login')
    title = request.form['title']
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("INSERT INTO books (title, class_id, is_published) VALUES (%s, %s, FALSE)", (title, class_id))
    conn.commit()
    conn.close()
    return redirect(f'/admin/books/{class_id}')

@app.route('/admin/books/publish/<int:book_id>', methods=['POST'])
def admin_books_publish(book_id):
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("UPDATE books SET is_published = TRUE WHERE id = %s", (book_id,))
    conn.commit()
    cursor.execute("SELECT class_id FROM books WHERE id = %s", (book_id,))
    class_id = cursor.fetchone()[0]
    cursor.close()
    conn.close()
    return redirect(f'/admin/books/{class_id}')

@app.route('/admin/mcqs/<int:book_id>')
def admin_mcqs(book_id):
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT 
            b.id,
            b.title,
            b.class_id,
            c.name AS class_name
        FROM books b
        JOIN classes c ON b.class_id = c.id
        WHERE b.id = %s
    """, (book_id,))
    book = cursor.fetchone()

    if not book:
        cursor.close()
        conn.close()
        return 'Course not found'

    cursor.execute("""
        SELECT 
            a.id,
            a.title,
            a.assessment_type,
            a.is_published,
            DATE_FORMAT(a.created_at, '%Y-%m-%d') AS created_date,
            COUNT(m.id) AS total_mcqs
        FROM assessments a
        LEFT JOIN mcqs m ON m.assessment_id = a.id
        WHERE a.book_id = %s
        GROUP BY a.id, a.title, a.assessment_type, a.is_published, a.created_at
        ORDER BY a.created_at DESC, a.id DESC
    """, (book_id,))
    assessments = cursor.fetchall()

    cursor.close()
    conn.close()

    return render_template(
        'admin/mcqs.html',
        book_title=book['title'],
        class_name=book['class_name'],
        book_id=book_id,
        class_id=book['class_id'],
        assessments=assessments
    )

@app.route('/admin/assessment/<int:assessment_id>/mcqs')
def admin_assessment_mcqs(assessment_id):
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT 
            a.id AS assessment_id,
            a.title AS assessment_title,
            a.is_published,
            DATE_FORMAT(a.created_at, '%Y-%m-%d') AS created_date,
            b.id AS book_id,
            b.title AS course_name,
            b.class_id,
            c.name AS class_name
        FROM assessments a
        JOIN books b ON a.book_id = b.id
        JOIN classes c ON b.class_id = c.id
        WHERE a.id = %s
    """, (assessment_id,))
    assessment = cursor.fetchone()

    if not assessment:
        cursor.close()
        conn.close()
        return "Assessment not found"

    cursor.execute("""
        SELECT *
        FROM mcqs
        WHERE assessment_id = %s
        ORDER BY id
    """, (assessment_id,))
    mcqs = cursor.fetchall()

    cursor.close()
    conn.close()

    return render_template(
        'admin/assessment_mcqs.html',
        assessment=assessment,
        mcqs=mcqs
    )

@app.route('/admin/mcqs/create/<int:book_id>', methods=['POST'])
def admin_mcqs_create(book_id):
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("""INSERT INTO mcqs 
        (question, option_a, option_b, option_c, option_d, correct_option, model_justification, book_id)
        VALUES (%s, %s, %s, %s, %s, %s, %s, %s)""",
        (request.form['question'], request.form['option_a'], request.form['option_b'],
         request.form['option_c'], request.form['option_d'], request.form['correct_option'],
         request.form['model_justification'], book_id))
    conn.commit()
    conn.close()
    return redirect(f'/admin/mcqs/{book_id}')

def extract_text_from_uploaded_file(uploaded_file):
    if not uploaded_file or not uploaded_file.filename:
        return ""

    filename = uploaded_file.filename.lower()
    suffix = os.path.splitext(filename)[1]

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        uploaded_file.save(tmp.name)
        temp_path = tmp.name

    try:
        if suffix in [".txt", ".md"]:
            with open(temp_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read().strip()

        if suffix == ".pdf":
            text_parts = []
            reader = PdfReader(temp_path)

            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

            return "\n".join(text_parts).strip()

        if suffix == ".docx":
            document = Document(temp_path)
            parts = []

            for paragraph in document.paragraphs:
                if paragraph.text.strip():
                    parts.append(paragraph.text.strip())

            for table in document.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))

            return "\n".join(parts).strip()

        if suffix == ".pptx":
            presentation = Presentation(temp_path)
            parts = []

            for slide_index, slide in enumerate(presentation.slides, start=1):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    parts.append(f"Slide {slide_index}:\n" + "\n".join(slide_text))

            return "\n".join(parts).strip()

        return ""

    finally:
        try:
            os.remove(temp_path)
        except Exception:
            pass
    
def select_relevant_text_for_topic(full_text, target_topic, max_chars=7000):
    """
    Select the most relevant portion of a full uploaded book/file
    for a requested chapter/topic before sending it to the LLM API.
    """
    import re

    full_text = (full_text or "").strip()
    target_topic = (target_topic or "").strip()

    if not full_text:
        return ""

    if not target_topic:
        return full_text[:max_chars]

    lower_text = full_text.lower()
    lower_topic = target_topic.lower()

    # 1. Exact phrase match
    exact_index = lower_text.find(lower_topic)
    if exact_index != -1:
        start = max(0, exact_index - max_chars // 3)
        end = min(len(full_text), start + max_chars)
        return full_text[start:end].strip()

    # 2. Search by topic keywords
    stopwords = {
        "the", "and", "for", "from", "with", "this", "that", "into",
        "about", "chapter", "topic", "unit", "lesson", "section",
        "of", "in", "on", "a", "an", "to", "is", "are"
    }

    topic_terms = [
        term for term in re.findall(r"[a-zA-Z0-9]+", lower_topic)
        if term not in stopwords and len(term) >= 2
    ]

    numbers = re.findall(r"\d+", lower_topic)
    topic_terms.extend(numbers)

    if not topic_terms:
        return full_text[:max_chars]

    words = full_text.split()
    chunk_size = 700
    overlap = 100

    chunks = []
    i = 0

    while i < len(words):
        chunk_words = words[i:i + chunk_size]
        chunk_text = " ".join(chunk_words)
        chunk_lower = chunk_text.lower()

        score = 0
        for term in topic_terms:
            score += chunk_lower.count(term)

        if lower_topic in chunk_lower:
            score += 20

        chunks.append({
            "index": len(chunks),
            "score": score,
            "text": chunk_text
        })

        i += chunk_size - overlap

    best_chunks = [chunk for chunk in chunks if chunk["score"] > 0]

    if not best_chunks:
        return ""

    best_chunks = sorted(best_chunks, key=lambda x: x["score"], reverse=True)[:4]
    best_chunks = sorted(best_chunks, key=lambda x: x["index"])

    selected_text = "\n\n".join(chunk["text"] for chunk in best_chunks).strip()

    return selected_text[:max_chars].strip()
    
@app.route('/admin/mcqs/generate/<int:book_id>', methods=['GET', 'POST'])
def admin_mcqs_generate(book_id):
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("SELECT title, class_id FROM books WHERE id = %s", (book_id,))
    book = cursor.fetchone()

    if not book:
        cursor.close()
        conn.close()
        return 'Book not found'

    generated_mcqs = None
    error = None
    generation_message = None

    notes_text = ''
    topic_name = ''
    target_topic = ''
    generation_mode = 'pasted_notes'
    num_mcqs = 3
    assessment_title = ''

    uploaded_filename = None
    extracted_file_chars = 0

    if request.method == 'POST':
        assessment_title = request.form.get('assessment_title', '').strip()
        notes_text = request.form.get('notes_text', '').strip()
        topic_name = request.form.get('topic_name', '').strip()
        target_topic = request.form.get('target_topic', '').strip()
        generation_mode = request.form.get('generation_mode', 'pasted_notes').strip()

        try:
            num_mcqs = int(request.form.get('num_mcqs', 3))
        except Exception:
            num_mcqs = 3

        num_mcqs = max(1, min(num_mcqs, 15))

        if not assessment_title:
            error = "Please enter a quiz or assessment name before generating MCQs."

        uploaded_file = request.files.get('notes_file')
        file_text = ""
        print("REQUEST FILES:", request.files)
        print("Uploaded file object:", uploaded_file)
        print("Uploaded filename:", uploaded_file.filename if uploaded_file else "NO FILE RECEIVED")

        if uploaded_file and uploaded_file.filename:
            uploaded_filename = uploaded_file.filename

            if not is_allowed_notes_file(uploaded_filename):
                error = "Unsupported file type. Please upload TXT, MD, PDF, DOCX, or PPTX."
            else:
                file_text = extract_text_from_uploaded_file(uploaded_file)
                extracted_file_chars = len(file_text)

                if not file_text:
                    error = (
                        "The uploaded file could not be prepared for MCQ generation. "
                        "Please try a smaller text-based file, paste the relevant notes, or add MCQs manually."
                )

        combined_notes_text = "\n\n".join(
            part for part in [notes_text, file_text] if part and part.strip()
        ).strip()

        if not error:
            if generation_mode == "topic_only":
                if not topic_name:
                    error = "Please enter a topic name for topic-only MCQ generation."

            elif generation_mode == "pasted_notes":
                if not combined_notes_text:
                    error = "Please paste notes text before generating MCQs."

            elif generation_mode == "uploaded_file":
                if not combined_notes_text:
                    error = "Please upload a supported file before generating MCQs."

            elif generation_mode == "uploaded_file_with_topic":
                if not combined_notes_text:
                    error = "Please upload a supported file before generating MCQs."
                elif not target_topic:
                    error = "Please enter the chapter/topic you want MCQs from."

            else:
                error = "Invalid generation mode selected."

        if not error:
            api_notes_text = combined_notes_text

            if generation_mode == "topic_only":
                api_notes_text = ""

            elif generation_mode == "uploaded_file_with_topic":
                api_notes_text = select_relevant_text_for_topic(
                    combined_notes_text,
                    target_topic,
                    max_chars=7000
                )

                print("Target topic:", target_topic)
                print("Relevant selected text length:", len(api_notes_text))
                print("Relevant selected text preview:", api_notes_text[:1000])

                if not api_notes_text.strip():
                    error = ( 
                        "The selected topic could not be prepared for generation. "
                        "Please use a clearer topic/chapter name, upload a smaller file, or paste the relevant section directly."
                    )

            else:
                # For pasted notes or full uploaded file, keep the input smaller for Colab.
                api_notes_text = combined_notes_text[:7000]

            print("Generation mode:", generation_mode)
            print("Original pasted notes length:", len(notes_text))
            print("Uploaded file extracted length:", len(file_text))
            print("Combined notes text length:", len(combined_notes_text))
            print("API notes text length:", len(api_notes_text))

            mcq_result = generate_mcqs_from_notes_api(
                notes_text=api_notes_text,
                num_mcqs=num_mcqs,
                topic_name=topic_name,
                target_topic=target_topic,
                generation_mode=generation_mode
            )

            print("MCQ API RESULT:", mcq_result)

            if isinstance(mcq_result, dict):
                generated_mcqs = mcq_result.get("mcqs", [])
                generation_message = mcq_result.get("message", "")
                generation_error = mcq_result.get("error")
            else:
                generated_mcqs = mcq_result or []
                generation_message = ""
                generation_error = None

            print("Generated MCQs count:", len(generated_mcqs))
            print("First generated MCQ:", generated_mcqs[0] if generated_mcqs else "NONE")

            requested_count = num_mcqs
            returned_count = len(generated_mcqs) if generated_mcqs else 0

            if not generated_mcqs:
                error = (
                    "MCQ generation could not be completed at this moment. "
                    "This may happen due to a weak internet connection, longer processing time, or temporary service interruption. "
                    "Please try again with fewer MCQs, shorter notes, or use manual MCQ entry to continue immediately."
                )

            elif returned_count < requested_count:
                generation_message = (
                    f"{returned_count} MCQ(s) were generated successfully. "
                    "The full requested count could not be completed at this moment. "
                    "You may review the generated MCQs, add more manually, or try again with fewer MCQs."
                )

            else:
                generation_message = (
                    f"{returned_count} MCQ(s) generated successfully. "
                    "Please review and edit them before saving the quiz."
                )

            

    cursor.close()
    conn.close()

    return render_template(
        'admin/generate_mcqs.html',
        book_id=book_id,
        book_title=book['title'],
        class_id=book['class_id'],
        generated_mcqs=generated_mcqs,
        error=error,
        generation_message=generation_message,
        assessment_title=assessment_title,
        notes_text=notes_text,
        topic_name=topic_name,
        target_topic=target_topic,
        generation_mode=generation_mode,
        num_mcqs=num_mcqs,
        uploaded_filename=uploaded_filename,
        extracted_file_chars=extracted_file_chars
    )

@app.route('/admin/mcqs/save-generated/<int:book_id>', methods=['POST'])
def admin_mcqs_save_generated(book_id):
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')

    assessment_title = request.form.get('assessment_title', '').strip()

    if not assessment_title:
        assessment_title = "Untitled Assessment"

    questions = request.form.getlist('question')
    option_as = request.form.getlist('option_a')
    option_bs = request.form.getlist('option_b')
    option_cs = request.form.getlist('option_c')
    option_ds = request.form.getlist('option_d')
    correct_options = request.form.getlist('correct_option')
    model_justifications = request.form.getlist('model_justification')

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute("""
        INSERT INTO assessments 
        (book_id, title, assessment_type, is_published, created_by, created_at)
        VALUES (%s, %s, 'quiz', 1, %s, NOW())
    """, (
        book_id,
        assessment_title,
        session.get('user_id')
    ))

    assessment_id = cursor.lastrowid

    saved_count = 0

    for i in range(len(questions)):
        question = questions[i].strip()
        option_a = option_as[i].strip()
        option_b = option_bs[i].strip()
        option_c = option_cs[i].strip()
        option_d = option_ds[i].strip()
        correct_option = correct_options[i].strip().upper()
        model_justification = model_justifications[i].strip()

        if correct_option not in ['A', 'B', 'C', 'D']:
            continue

        if not question or not option_a or not option_b or not option_c or not option_d or not model_justification:
            continue

        cursor.execute("""
            INSERT INTO mcqs 
            (
                question,
                option_a,
                option_b,
                option_c,
                option_d,
                correct_option,
                model_justification,
                book_id,
                assessment_id
            )
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
        """, (
            question,
            option_a,
            option_b,
            option_c,
            option_d,
            correct_option,
            model_justification,
            book_id,
            assessment_id
        ))

        saved_count += 1

    if saved_count == 0:
        cursor.execute("DELETE FROM assessments WHERE id = %s", (assessment_id,))
        conn.commit()
        cursor.close()
        conn.close()
        return redirect(f'/admin/mcqs/generate/{book_id}')

    conn.commit()
    cursor.close()
    conn.close()

    return redirect(f'/admin/assessment/{assessment_id}/mcqs')

@app.route('/admin/students/<int:class_id>')
def admin_students(class_id):
    if session.get('user_role') != 'administrator':
        return redirect('/login')
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    cursor.execute("SELECT name FROM classes WHERE id = %s AND created_by = %s", (class_id, session['user_id']))
    cls = cursor.fetchone()
    if not cls:
        conn.close()
        return 'Class not found'
    cursor.execute("SELECT * FROM students WHERE class_id = %s", (class_id,))
    students = cursor.fetchall()
    conn.close()
    return render_template('admin/students.html', class_name=cls['name'], students=students, class_id=class_id)

@app.route('/admin/students/add/<int:class_id>', methods=['POST'])
def admin_students_add(class_id):
    if session.get('user_role') != 'administrator':
        return redirect('/login')
    name = request.form['name']
    roll = request.form['roll_number']
    username = request.form['username']
    password = request.form['password']
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("INSERT INTO students (name, roll_number, class_id, username, password) VALUES (%s, %s, %s, %s, %s)",
                   (name, roll, class_id, username, password))
    conn.commit()
    conn.close()
    return redirect(f'/admin/students/{class_id}')

@app.route('/admin/mcqs/delete/<int:mcq_id>', methods=['POST'])
def admin_mcqs_delete(mcq_id):
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("DELETE FROM results WHERE mcq_id = %s", (mcq_id,))
    cursor.execute("DELETE FROM mcqs WHERE id = %s", (mcq_id,))
    conn.commit()
    cursor.close()
    conn.close()
    return redirect(request.referrer or '/admin/classes')

@app.route('/admin/results')
def admin_results():
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    cursor.execute("SELECT id, name FROM classes")
    classes = cursor.fetchall()
    results = {}
    for cls in classes:
        cursor.execute("""
            SELECT s.id as student_id, s.name as student_name, s.roll_number, b.title as book_title, b.id as book_id,
                   ROUND(SUM(r.marks) / COUNT(r.id) * 100, 2) AS total_percentage
            FROM results r
            JOIN mcqs m ON r.mcq_id = m.id
            JOIN books b ON m.book_id = b.id
            JOIN students s ON r.student_id = s.id
            WHERE s.class_id = %s
            GROUP BY s.id, b.id
            ORDER BY s.name, b.title
        """, (cls['id'],))
        results[cls['id']] = cursor.fetchall()
    conn.close()
    return render_template('admin/results.html', classes=classes, results=results)

@app.route('/admin/results/student/<int:student_id>/<int:book_id>')
def admin_results_detail(student_id, book_id):
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    cursor.execute("SELECT name FROM students WHERE id = %s", (student_id,))
    student = cursor.fetchone()
    cursor.execute("SELECT title FROM books WHERE id = %s", (book_id,))
    book = cursor.fetchone()
    cursor.execute("""
        SELECT m.question, r.selected_option, r.marks, r.feedback
        FROM results r
        JOIN mcqs m ON r.mcq_id = m.id
        WHERE r.student_id = %s AND m.book_id = %s
        ORDER BY m.id
    """, (student_id, book_id))
    details = cursor.fetchall()
    conn.close()
    return render_template('admin/result_detail.html', student=student['name'], book=book['title'], details=details)

@app.route('/admin/teachers')
def admin_teachers():
    if session.get('user_role') != 'administrator':
        return redirect('/login')
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    cursor.execute("SELECT id, username, name FROM admins WHERE role = 'teacher'")
    teachers = cursor.fetchall()
    conn.close()
    return render_template('admin/teachers.html', teachers=teachers)

@app.route('/admin/teachers/add', methods=['POST'])
def admin_teachers_add():
    if session.get('user_role') != 'administrator':
        return redirect('/login')
    name = request.form['name']
    username = request.form['username']
    password = request.form['password']
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("INSERT INTO admins (username, password, role, name) VALUES (%s, %s, 'teacher', %s)", (username, password, name))
    conn.commit()
    conn.close()
    return redirect('/admin/teachers')

@app.route('/student/results/detail/<int:book_id>')
def student_results_detail(book_id):
    if session.get('user_role') != 'student':
        return redirect('/login')

    student_id = session['user_id']

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT 
            s.name,
            s.roll_number,
            c.name AS class_name
        FROM students s
        JOIN classes c ON s.class_id = c.id
        WHERE s.id = %s
    """, (student_id,))
    student = cursor.fetchone()

    cursor.execute("""
        SELECT 
            b.id,
            b.title,
            c.name AS class_name
        FROM books b
        JOIN classes c ON b.class_id = c.id
        JOIN students s ON s.class_id = c.id
        WHERE b.id = %s
          AND s.id = %s
    """, (book_id, student_id))
    course = cursor.fetchone()

    if not course:
        cursor.close()
        conn.close()
        return "Course not found"

    cursor.execute("""
        SELECT 
            a.id AS assessment_id,
            a.title AS assessment_title,
            DATE_FORMAT(MIN(r.submitted_at), '%Y-%m-%d') AS attempt_date,
            ROUND(SUM(r.marks), 2) AS obtained_marks,
            COUNT(r.id) AS total_questions,
            ROUND((SUM(r.marks) / COUNT(r.id)) * 100, 2) AS percentage
        FROM results r
        JOIN assessments a ON r.assessment_id = a.id
        WHERE r.student_id = %s
          AND a.book_id = %s
        GROUP BY a.id, a.title
        ORDER BY MIN(r.submitted_at) DESC, a.id DESC
    """, (student_id, book_id))
    assessments = cursor.fetchall()

    for assessment in assessments:
        cursor.execute("""
            SELECT 
                m.question,
                m.option_a,
                m.option_b,
                m.option_c,
                m.option_d,
                m.correct_option,
                m.model_justification,
                r.selected_option,
                r.justification AS student_justification,
                ROUND(r.marks, 2) AS marks,
                r.feedback,
                r.llm_feedback
            FROM results r
            JOIN mcqs m ON r.mcq_id = m.id
            WHERE r.student_id = %s
              AND r.assessment_id = %s
            ORDER BY m.id
        """, (student_id, assessment["assessment_id"]))

        assessment["details"] = cursor.fetchall()

    cursor.close()
    conn.close()

    return render_template(
        'student/result_detail.html',
        student=student,
        course=course,
        assessments=assessments
    )

@app.route('/admin/profile')
def admin_profile():
    if session.get('user_role') not in ['administrator', 'teacher']:
        return redirect('/login')

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT id, name, username, role
        FROM admins
        WHERE id = %s
    """, (session['user_id'],))

    user = cursor.fetchone()

    cursor.close()
    conn.close()

    return render_template('admin/profile.html', user=user or {})

@app.route('/admin/students')
def admin_students_overview():
    if session.get('user_role') != 'administrator':
        return redirect('/login')

    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    cursor.execute("""
        SELECT id, name
        FROM classes
        ORDER BY name
    """)
    classes = cursor.fetchall()

    cursor.execute("""
        SELECT 
            s.id,
            s.name,
            s.roll_number,
            s.username,
            c.name AS class_name,
            c.id AS class_id
        FROM students s
        LEFT JOIN classes c ON s.class_id = c.id
        ORDER BY c.name, s.name
    """)
    students = cursor.fetchall()

    cursor.close()
    conn.close()

    return render_template(
        'admin/students_overview.html',
        classes=classes,
        students=students
    )


@app.route('/admin/students/add', methods=['POST'])
def admin_students_add_overview():
    if session.get('user_role') != 'administrator':
        return redirect('/login')

    class_id = request.form.get('class_id')
    name = request.form.get('name', '').strip()
    roll_number = request.form.get('roll_number', '').strip()
    username = request.form.get('username', '').strip()
    password = request.form.get('password', '').strip()

    if not class_id or not name or not username or not password:
        return redirect('/admin/students')

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute("""
        INSERT INTO students (name, roll_number, class_id, username, password)
        VALUES (%s, %s, %s, %s, %s)
    """, (name, roll_number, class_id, username, password))

    conn.commit()
    cursor.close()
    conn.close()

    return redirect('/admin/students')


@app.route('/admin/students/delete/<int:student_id>', methods=['POST'])
def admin_students_delete(student_id):
    if session.get('user_role') != 'administrator':
        return redirect('/login')

    conn = get_db_connection()
    cursor = conn.cursor()

    try:
        cursor.execute("DELETE FROM results WHERE student_id = %s", (student_id,))

        try:
            cursor.execute("DELETE FROM student_assessment_attempts WHERE student_id = %s", (student_id,))
        except Exception:
            pass

        cursor.execute("DELETE FROM students WHERE id = %s", (student_id,))
        conn.commit()

    except Exception as e:
        print("Student delete failed:", e)
        conn.rollback()

    cursor.close()
    conn.close()

    return redirect('/admin/students')

if __name__ == '__main__':
    app.run(debug=True)